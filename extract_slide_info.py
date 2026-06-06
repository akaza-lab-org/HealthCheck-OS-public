import os
import pdfplumber
from pptx import Presentation
import json

def get_pdf_titles(path):
    titles = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    first_line = text.split('\n')[0].strip()
                    titles.append({"page": i + 1, "title": first_line})
                else:
                    titles.append({"page": i + 1, "title": "[Image/No Text]"})
    except Exception as e:
        return str(e)
    return titles

def get_pptx_titles(path):
    titles = []
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            title = "[No Title]"
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
            elif len(slide.shapes) > 0:
                # Try to find the largest text box or first text frame
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        title = shape.text.strip().split('\n')[0]
                        break
            titles.append({"page": i + 1, "title": title})
    except Exception as e:
        return str(e)
    return titles

# Target files
files = {
    "wegovy_env": r"G:\マイドライブ\赤座至先生_ウゴービスライドセット\赤座至先生_ウゴービスライドセット\AID9568】肥満・肥満症を取り巻く環境と、医学的介入の意義（スライドセット・医師向け）JP25OB00198v2.pptx",
    "wegovy_select": r"G:\マイドライブ\赤座至先生_ウゴービスライドセット\赤座至先生_ウゴービスライドセット\【AID9643】医師提供用_SELECT 肥満指標別サブグループ解析(16 9)_JP26SEMO00080.pptx",
    "wegovy_story": r"G:\マイドライブ\赤座至先生_ウゴービスライドセット\赤座至先生_ウゴービスライドセット\【AID9643】医師提供用_ウゴービブランドストーリー_スピーカースライド (16：9)_JP26SEMO00041.pptx",
    "wegovy_actionio": r"G:\マイドライブ\赤座至先生_ウゴービスライドセット\赤座至先生_ウゴービスライドセット\医師提供用ActionIOスライド_JP25CO00010 (1).pptx",
    "wegovy_metabolism": r"G:\マイドライブ\赤座至先生_ウゴービスライドセット\赤座至先生_ウゴービスライドセット\医師提供用_肥満・肥満症における代謝適応とGLP-1受容体作動薬_16_9_JP24SEMO00184.pptx",
    "wegovy_guideline": r"G:\マイドライブ\赤座至先生_ウゴービスライドセット\赤座至先生_ウゴービスライドセット\医師提供用スライド_肥満症診療ガイドラインスライド_16_9_JP24SEMO00126.pptx",
    "past_edu": r"G:\マイドライブ\糖尿病講義スライド\教育入院\【最終版】202306糖尿病治療の進め方と教育入院の適応40min.pptx",
    "past_glp1_diff": r"G:\マイドライブ\糖尿病講義スライド\治療\GLP1RA薬の製剤間の違い_22022.3.pptx"
}

results = {}
for key, path in files.items():
    if path.endswith(".pptx"):
        results[key] = get_pptx_titles(path)
    elif path.endswith(".pdf"):
        results[key] = get_pdf_titles(path)

print(json.dumps(results, ensure_ascii=False, indent=2))
